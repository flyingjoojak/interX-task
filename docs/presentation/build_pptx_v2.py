"""InterX 최종 면접 발표 v2 — 19장 .pptx 생성기.

사용법:
    python docs/presentation/build_pptx_v2.py

출력:
    docs/presentation/interX_발표_v2.pptx
"""
from __future__ import annotations

import os
import sqlite3
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
OUT = DOCS / "presentation" / "interX_발표_v2.pptx"
DB = ROOT / "backend" / "interx.db"

PRIMARY = RGBColor(0xFF, 0x80, 0x00)
INK = RGBColor(0x11, 0x18, 0x27)
SUBINK = RGBColor(0x37, 0x41, 0x51)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "맑은 고딕"
TOTAL = 19


# ---------- 기본 헬퍼 ----------

def set_run(run, *, size=14, color=INK, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color, bold=bold)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=BORDER, line_width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_picture_fit(slide, img_path, left, top, box_w, box_h):
    """이미지를 (left, top, box_w, box_h) 안에 종횡비 유지하며 중앙 배치."""
    p = str(img_path)
    if not os.path.exists(p):
        return None
    with PILImage.open(p) as im:
        iw, ih = im.size
    if iw == 0 or ih == 0:
        return None
    ir = iw / ih
    br = box_w / box_h
    if ir > br:
        w = int(box_w)
        h = int(box_w / ir)
    else:
        h = int(box_h)
        w = int(box_h * ir)
    x = int(left + (box_w - w) / 2)
    y = int(top + (box_h - h) / 2)
    return slide.shapes.add_picture(p, x, y, width=w, height=h)


def add_footer(slide, idx):
    add_textbox(slide, Inches(0.4), Inches(7.05), Inches(6), Inches(0.3),
                "주재홍 · InterX 2026", size=10, color=MUTED)
    add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
                f"{idx} / {TOTAL}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                title, size=28, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(1.05), Inches(12), Inches(0.45),
                    subtitle, size=13, color=MUTED)
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(1.0), Emu(36000),
             fill=PRIMARY, line=PRIMARY)


def add_card(slide, left, top, width, height, *,
             accent_top=True, accent=PRIMARY):
    add_rect(slide, left, top, width, height, fill=WHITE, line=BORDER)
    if accent_top:
        add_rect(slide, left, top, width, Inches(0.08),
                 fill=accent, line=accent)


def add_bullets(slide, left, top, width, height, bullets, *,
                size=12, color=SUBINK, bullet_color=PRIMARY,
                line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = "•  "
        set_run(r1, size=size, color=bullet_color, bold=True)
        r2 = p.add_run()
        r2.text = b
        set_run(r2, size=size, color=color)
        p.space_after = Pt(3)
    return tb


def add_arrow_step(slide, left, top, width, height, label, *,
                   accent=False):
    fill = PRIMARY if accent else WHITE
    border = PRIMARY if accent else BORDER
    text_color = WHITE if accent else INK
    add_rect(slide, left, top, width, height, fill=fill, line=border)
    add_textbox(slide, left, top, width, height, label,
                size=11, bold=accent, color=text_color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_chevron(slide, left, top, height=Inches(0.18)):
    add_textbox(slide, left, top, Inches(0.3), height, "▶",
                size=11, color=PRIMARY, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def fetch_step_costs():
    """backend/interx.db 에서 step별 실측 비용을 읽어 정렬된 리스트로."""
    if not DB.exists():
        return []
    con = sqlite3.connect(str(DB))
    cur = con.cursor()
    cur.execute(
        "SELECT step, SUM(cost_usd), COUNT(*) "
        "FROM token_usage GROUP BY step ORDER BY 2 DESC"
    )
    rows = cur.fetchall()
    con.close()
    return [(step, float(cost), int(n)) for step, cost, n in rows if cost]


# ---------- 슬라이드 ----------

def slide_1(prs):
    """타이틀."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(9.6), Inches(0), Inches(3.83), Inches(7.5),
             fill=PRIMARY, line=PRIMARY)
    add_rect(slide, Inches(10.4), Inches(5.6), Inches(0.9), Inches(0.9),
             fill=INK, line=INK)
    add_rect(slide, Inches(11.6), Inches(5.6), Inches(0.9), Inches(0.9),
             fill=WHITE, line=WHITE)

    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(8.5), Inches(0.7),
                "InterX 합류를 위한", size=22, color=MUTED)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(8.5), Inches(1.3),
                "과제 발표", size=58, bold=True, color=INK)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(9), Inches(0.6),
                '"단순 구현보다 0.1%의 효율을 중시하는 AI Engineer"',
                size=18, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.7), Inches(6.3), Inches(8), Inches(0.4),
                "주재홍   ·   2026.04.23", size=13, color=MUTED)
    add_textbox(slide, Inches(0.7), Inches(6.7), Inches(8), Inches(0.4),
                "kinduki123@gmail.com   ·   github.com/flyingjoojak",
                size=11, color=MUTED)


def slide_2(prs):
    """저는 이런 엔지니어입니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "저는 이런 엔지니어입니다",
              "10년+ 개발 경험과 컴퓨터공학 전공의 결합")

    add_textbox(slide, Inches(0.6), Inches(1.85), Inches(12.2), Inches(0.7),
                '"돌아가는 AI"가 아니라 "효율로 증명되는 AI"를 만듭니다.',
                size=22, bold=True, color=INK)

    cards = [
        ("리소스 최적화",
         "ONNX · 양자화 · 캐싱으로\n하드웨어와 API 비용을 함께 절감",
         "RapidOCR 70%↓  ·  SDXL 90%↓",
         PRIMARY),
        ("고신뢰 RAG",
         "Hybrid Search + Reranking 으로\n토큰을 절약하며 정확도를 상승",
         "관련성 +45%  ·  토큰 −30%",
         INK),
        ("엔드-투-엔드 책임감",
         "기획·모델·서빙·프론트·운영까지\n직접 만들어 본 5개 프로젝트",
         "비전 · LLM · 풀스택 · DevOps",
         INK),
    ]
    card_w = Inches(4.0)
    card_h = Inches(3.6)
    gap = Inches(0.18)
    start_left = Inches(0.6)
    top = Inches(2.95)
    for i, (title, body, metric, accent) in enumerate(cards):
        left = start_left + (card_w + gap) * i
        add_card(slide, left, top, card_w, card_h, accent=accent)
        add_textbox(slide, left + Inches(0.3), top + Inches(0.35),
                    card_w - Inches(0.6), Inches(0.55),
                    title, size=18, bold=True, color=accent)
        add_textbox(slide, left + Inches(0.3), top + Inches(1.05),
                    card_w - Inches(0.6), Inches(1.5),
                    body, size=13, color=SUBINK)
        add_rect(slide, left + Inches(0.3), top + card_h - Inches(0.85),
                 card_w - Inches(0.6), Inches(0.55),
                 fill=LIGHT, line=LIGHT)
        add_textbox(slide, left + Inches(0.3), top + card_h - Inches(0.85),
                    card_w - Inches(0.6), Inches(0.55),
                    metric, size=12, bold=True, color=INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 2)


def slide_3(prs):
    """커리어 하이라이트 — 5개 프로젝트 카드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "커리어 하이라이트",
              "엔드-투-엔드로 만든 5개 프로젝트 — 정량 지표로 증명")

    items = [
        ("AI 광고 생성 파이프라인",
         "소상공인 마케팅 자동화\n로컬 추론으로 비용 0원",
         "90%↓ / 0원",
         "300s → 30s · API 비용 0원",
         "SDXL · Diffusers · FastAPI · FP16/INT8"),
        ("RFP RAG 분석 시스템",
         "입찰 제안서 자동 분석\n근거 인용으로 환각 억제",
         "+45% / −30%",
         "관련성 +45% · 토큰 −30%",
         "LangChain · ChromaDB · FlashRank · Ragas"),
        ("약제 식별 시스템",
         "조제 오류 방지를 위한\n비전 검수 자동화",
         "99.1%",
         "750회 반복 실험",
         "YOLOv11 · Optuna · Ensemble · OpenCV"),
        ("디스코드 커뮤니티 플랫폼",
         "활성도 기반 랭킹\n가상 재화 경제 시스템",
         "₩490K",
         "100+ 서버 · DAU 300",
         "Next.js · FastAPI · Pycord · PG · Redis"),
        ("실시간 번역 오버레이",
         "외국어 콘텐츠 실시간 번역\n저사양에서도 상시 구동",
         "70%↓ / 0원",
         "OCR 70%↓ · 캐시로 비용 0원",
         "PyQt6 · RapidOCR · ONNX Runtime"),
    ]

    card_w = Inches(2.45)
    card_h = Inches(4.85)
    gap = Inches(0.12)
    total_w = card_w * 5 + gap * 4
    start_left = (Inches(13.333) - total_w) / 2
    top = Inches(1.95)
    for i, (title, sub, big, metric, stack) in enumerate(items):
        left = start_left + (card_w + gap) * i
        add_card(slide, left, top, card_w, card_h, accent=PRIMARY)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                    card_w - Inches(0.4), Inches(0.4),
                    f"0{i+1}", size=12, bold=True, color=PRIMARY)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.55),
                    card_w - Inches(0.4), Inches(0.9),
                    title, size=13, bold=True, color=INK,
                    line_spacing=1.15)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.5),
                    card_w - Inches(0.4), Inches(0.95),
                    sub, size=10, color=MUTED, line_spacing=1.2)
        add_textbox(slide, left + Inches(0.2), top + Inches(2.55),
                    card_w - Inches(0.4), Inches(0.7),
                    big, size=22, bold=True, color=PRIMARY,
                    align=PP_ALIGN.CENTER)
        add_rect(slide, left + Inches(0.2), top + Inches(3.35),
                 card_w - Inches(0.4), Inches(0.5),
                 fill=LIGHT, line=LIGHT)
        add_textbox(slide, left + Inches(0.2), top + Inches(3.35),
                    card_w - Inches(0.4), Inches(0.5),
                    metric, size=10, color=INK, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left + Inches(0.2), top + Inches(4.0),
                    card_w - Inches(0.4), Inches(0.75),
                    stack, size=9, color=MUTED, line_spacing=1.25)

    add_textbox(slide, Inches(0.6), Inches(7.0), Inches(12.2), Inches(0.35),
                "왼쪽 4개 = 효율(얼마나 줄였나)  ·  오른쪽 1개 = 실전성(얼마나 빨리 통했나)",
                size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)


def slide_4(prs):
    """지원 동기."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "지원 동기 — 왜 인터엑스인가",
              "AI를 써서 무엇을 바꾸는가에 답이 있는 회사")

    points = [
        ("AI 에이전트로 채용을 재정의",
         "단순 ATS 자동화가 아니라, 사람의 판단을 정량적 근거와 이중 검증으로\n"
         "보조하는 방향. 제가 추구해온 \"효율로 증명되는 AI\"와 정확히 같은 지향."),
        ("정형화되지 않은 어려운 문제",
         "이력서 환각, 면접 대화의 모순 탐지, 설명 가능한 AI 결과 — 전부 RAG에서\n"
         "토큰·정확도 트레이드오프를 직접 설계해 본 제 경험이 그대로 쓰일 영역."),
        ("AX(AI Experience) 관점의 프로덕트",
         "모델 성능 그 자체보다 면접관·지원자에게 실제 가치로 연결되는 지점을 설계.\n"
         "디스코드 봇에서 DAU 300 · 첫주 49만원으로 검증해 본 \"현장 적합도\"와 직결."),
    ]

    top = Inches(2.0)
    row_h = Inches(1.55)
    gap = Inches(0.15)
    for i, (head, body) in enumerate(points):
        y = top + (row_h + gap) * i
        add_rect(slide, Inches(0.6), y, Inches(12.2), row_h,
                 fill=WHITE, line=BORDER)
        add_rect(slide, Inches(0.6), y, Inches(0.14), row_h,
                 fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, Inches(0.95), y + Inches(0.3),
                    Inches(0.6), Inches(0.5),
                    f"0{i+1}", size=20, bold=True, color=PRIMARY)
        add_textbox(slide, Inches(1.7), y + Inches(0.18),
                    Inches(11), Inches(0.5),
                    head, size=17, bold=True, color=INK)
        add_textbox(slide, Inches(1.7), y + Inches(0.65),
                    Inches(11), Inches(0.95),
                    body, size=12, color=SUBINK, line_spacing=1.3)

    add_footer(slide, 4)


def slide_5(prs):
    """12가치 매칭 + 입사 후 포부."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "12가치 매칭 + 입사 후 포부",
              "12가치 중 4개에 먼저 기여하겠습니다")

    values = [
        ("끈기",
         "약제 식별 프로젝트에서 750회 이상의 반복 실험과 정밀 튜닝을 거쳐\n"
         "인식 정확도 99.1% 기록."),
        ("지속적개선",
         "RFP RAG에 FlashRank 리랭킹 도입으로 노이즈 제거 + 연관성 높은 컨텍스트 선별.\n"
         "MobileNetv2 기반 RapidOCR을 ONNX로 최적화해 부하 최소화 및 추론 구현.\n"
         "양자화·어텐션 최적화로 하드웨어 추론 성능 확보."),
        ("문제해결",
         "조제실 내 오투약 리스크를 줄이기 위한 약제 식별 시스템 구축.\n"
         "외국어 게임/영상 시청 시 발생하는 정보 격차를 실시간 번역 오버레이로 해소.\n"
         "클릭 간격 표준편차 기반 행동 패턴 알고리즘으로 매크로 방지 로직 적용."),
        ("비판적사고",
         "단순 LLM 호출에서 벗어나 Hybrid Search + FlashRank 리랭킹으로\n"
         "토큰 낭비 절감과 신뢰도 확보. 운영비 절감을 위해 유료 API 대신\n"
         "로컬 허깅페이스(SDXL) 모델 및 서빙 환경 구축. YOLOv11과 Faster R-CNN을\n"
         "비교하고 Optuna로 하이퍼파라미터 최적화 수행."),
    ]
    top = Inches(1.95)
    row_h = Inches(1.05)
    gap = Inches(0.08)
    for i, (name, body) in enumerate(values):
        y = top + (row_h + gap) * i
        add_rect(slide, Inches(0.6), y, Inches(12.2), row_h,
                 fill=WHITE, line=BORDER)
        add_rect(slide, Inches(0.6), y, Inches(2.4), row_h,
                 fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, Inches(0.6), y, Inches(2.4), row_h,
                    name, size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.15), y + Inches(0.1),
                    Inches(9.5), row_h - Inches(0.2),
                    body, size=11, color=SUBINK, line_spacing=1.3)

    pledge_top = Inches(6.6)
    add_rect(slide, Inches(0.6), pledge_top, Inches(12.2), Inches(0.55),
             fill=LIGHT, line=LIGHT)
    add_textbox(slide, Inches(0.6), pledge_top, Inches(12.2), Inches(0.55),
                "포부 — 회귀 테스트 · 비용 계측 · 자가 검증 루프로 네 가치를 코드 레이어에 박아 넣겠습니다.",
                size=12, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 5)


def slide_6(prs):
    """과제 오버뷰 — 9단계."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "과제 오버뷰 — InterX MVP",
              "FastAPI · LangGraph · Next.js 풀스택")

    # 좌측: 홈 스크린샷
    home = DOCS / "home.png"
    sx, sy = Inches(0.6), Inches(1.95)
    sw, sh = Inches(6.4), Inches(4.6)
    add_rect(slide, sx, sy, sw, sh, fill=LIGHT, line=BORDER)
    add_picture_fit(slide, home, sx + Inches(0.05), sy + Inches(0.05),
                    sw - Inches(0.1), sh - Inches(0.1))

    # 우측: 플로우 (9단계)
    fx = Inches(7.25)
    fw = Inches(5.55)
    fy = Inches(1.95)
    add_textbox(slide, fx, fy, fw, Inches(0.4),
                "데이터 흐름 (9단계)", size=14, bold=True, color=PRIMARY)

    steps = [
        "① 이력서/포트폴리오 업로드 (PDF·JPG·PNG)",
        "② Claude Vision OCR — 읽기 순서 강제",
        "③ 구조화 추출 (6필드 JSON)",
        "④ 12가치 매핑 (score·evidence·examples)",
        "⑤ 자가검증 루프 (substring · token-overlap)",
        "⑥ 문서 신뢰도 산출",
        "⑦ 모순 탐지 (4축 · severity)",
        "⑧ 사전 압박 질문 생성",
        "⑨ 실시간 면접 (꼬리질문) → PDF 리포트",
    ]
    step_h = Inches(0.38)
    step_gap = Inches(0.05)
    sy2 = fy + Inches(0.5)
    for i, s in enumerate(steps):
        y = sy2 + (step_h + step_gap) * i
        add_rect(slide, fx, y, fw, step_h, fill=WHITE, line=BORDER)
        accent_fill = PRIMARY if i in (3, 4) else PRIMARY  # emphasize 12값+자가검증
        add_rect(slide, fx, y, Inches(0.1), step_h,
                 fill=accent_fill, line=accent_fill)
        add_textbox(slide, fx + Inches(0.2), y, fw - Inches(0.3), step_h,
                    s, size=10, color=INK,
                    anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.35),
                "asyncio.Queue 단일 워커로 LLM 호출 직렬화  ·  모든 호출은 usage_scope로 phase/step 자동 태깅",
                size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def slide_7(prs):
    """미션 1 [수행]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 1. 데이터 파이프라인  ·  [수행]",
              "비정형 PDF/이미지 → 구조화 JSON → 12가치 매핑")

    add_textbox(slide, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.4),
                "파이프라인 4단계", size=14, bold=True, color=PRIMARY)
    stages = [
        ("01  OCR",
         "Claude Vision 기반.  \"위→아래, 좌→우, 2-컬럼 블록 유지\"\n읽기 순서를 프롬프트로 강제."),
        ("02  익명화 (PII-safe)",
         "이메일·전화번호를 토큰 치환.\n분석 후 복원 — 프롬프트 오염 차단."),
        ("03  구조화 추출",
         "6필드 JSON: career / education / skills /\nachievements / certifications / projects."),
        ("04  12가치 매핑",
         "가치별 score(0–100) · evidence(1–2문장) ·\nexamples(원문 인용 배열) 생성."),
    ]
    sy = Inches(2.4)
    sh = Inches(1.05)
    sgap = Inches(0.1)
    for i, (head, body) in enumerate(stages):
        y = sy + (sh + sgap) * i
        add_rect(slide, Inches(0.6), y, Inches(6.0), sh,
                 fill=WHITE, line=BORDER)
        add_rect(slide, Inches(0.6), y, Inches(0.12), sh,
                 fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, Inches(0.85), y + Inches(0.13),
                    Inches(5.6), Inches(0.4),
                    head, size=14, bold=True, color=INK)
        add_textbox(slide, Inches(0.85), y + Inches(0.5),
                    Inches(5.6), Inches(0.55),
                    body, size=11, color=SUBINK, line_spacing=1.25)

    rx = Inches(7.0)
    rw = Inches(5.85)
    add_textbox(slide, rx, Inches(1.95), rw, Inches(0.4),
                "프롬프트 엔지니어링 디테일", size=14, bold=True, color=PRIMARY)
    add_rect(slide, rx, Inches(2.4), rw, Inches(4.55),
             fill=LIGHT, line=BORDER)
    bullets = [
        "부정 예시 명시 — \"열정적입니다\", \"책임감이 강합니다\"\n     같은 추상적 자기서술은 증거로 인정하지 않음.",
        "80점 하한 임계선 — 충분한 구체적 증거(수치·프로젝트·결과)가\n     반복 관찰될 때만 80↑ (100점까지 열려 있음).",
        "거절 기본값 — 증거가 없는 가치는 0–30 부여하고\n     examples는 빈 배열.",
        "원문 인용 강제 — examples는 이력서의 구체 문장을\n     원문 그대로 1–3개 인용.",
    ]
    add_bullets(slide, rx + Inches(0.25), Inches(2.6),
                rw - Inches(0.5), Inches(4.2),
                bullets, size=12, color=SUBINK, line_spacing=1.4)

    add_footer(slide, 7)


def slide_8(prs):
    """미션 1 [결과]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 1. 데이터 파이프라인  ·  [결과]",
              "구조화 데이터셋 + 12가치 점수 + 원문 증거가 한 객체로")

    lx = Inches(0.6)
    lw = Inches(6.6)
    add_textbox(slide, lx, Inches(1.95), lw, Inches(0.4),
                "결과물", size=14, bold=True, color=PRIMARY)

    blocks = [
        ("DB 스키마",
         "Analysis.structured_data (6필드 JSON)\nAnalysis.values_scores (12가치별 점수+근거)"),
        ("API",
         "GET /api/candidates/{id}/analysis\nGET /api/candidates/{id}/debug/raw"),
        ("프론트 컴포넌트",
         "ValueRadarChart · ValueList · EvidencePanel · RawDataViewer"),
    ]
    by = Inches(2.4)
    bh = Inches(1.0)
    bgap = Inches(0.12)
    for i, (h, b) in enumerate(blocks):
        y = by + (bh + bgap) * i
        add_rect(slide, lx, y, lw, bh, fill=WHITE, line=BORDER)
        add_rect(slide, lx, y, Inches(0.1), bh, fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, lx + Inches(0.25), y + Inches(0.12),
                    lw - Inches(0.4), Inches(0.4),
                    h, size=13, bold=True, color=INK)
        add_textbox(slide, lx + Inches(0.25), y + Inches(0.48),
                    lw - Inches(0.4), bh - Inches(0.55),
                    b, size=11, color=SUBINK, line_spacing=1.3)

    jy = Inches(5.85)
    add_rect(slide, lx, jy, lw, Inches(1.1), fill=INK, line=INK)
    add_textbox(slide, lx + Inches(0.2), jy + Inches(0.08),
                lw - Inches(0.4), Inches(0.95),
                'values_scores.지속적개선 = {\n  "score": 88,\n  "examples": ["ONNX 가속으로 추론 70% 단축", ...]\n}',
                size=10, color=WHITE)

    rx = Inches(7.4)
    ry = Inches(1.95)
    rw = Inches(5.45)
    rh = Inches(5.0)
    add_rect(slide, rx, ry, rw, rh, fill=LIGHT, line=BORDER)
    add_picture_fit(slide, DOCS / "report.png",
                    rx + Inches(0.05), ry + Inches(0.05),
                    rw - Inches(0.1), rh - Inches(0.1))
    add_textbox(slide, rx, ry + rh + Inches(0.02), rw, Inches(0.3),
                "리포트 화면 — 12가치 레이더 + 원문 증거",
                size=10, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 8)


def slide_9(prs):
    """미션 2 [수행] — LangGraph 두 개."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 2. 모순 탐지 + 꼬리질문  ·  [수행]",
              "LangGraph 두 개 — 사전 분석 그래프 + 실시간 면접 그래프")

    add_textbox(slide, Inches(0.6), Inches(1.95), Inches(12.2), Inches(0.4),
                "체인 ①  분석 그래프 (사전, analysis_graph.py)",
                size=13, bold=True, color=PRIMARY)

    chain1 = [
        "parse_documents", "anonymize_pii", "extract_structured",
        "score_12_values", "self_verify_evidence",
        "doc_reliability", "detect_contradictions",
        "preemptive_questions", "compile_restore",
    ]
    cw = Inches(1.32)
    ch = Inches(0.55)
    cgap = Inches(0.04)
    cy = Inches(2.4)
    cx = Inches(0.6)
    for i, name in enumerate(chain1):
        x = cx + (cw + cgap) * i
        accent = name in ("self_verify_evidence", "detect_contradictions")
        add_arrow_step(slide, x, cy, cw, ch, name, accent=accent)

    add_textbox(slide, Inches(0.6), Inches(3.05), Inches(12.2), Inches(0.4),
                "모순 탐지 4축: 날짜/기간 · 직책/역할 · 성과 수치 · 기술 스택  ·  severity high/medium/low",
                size=11, color=MUTED)
    add_textbox(slide, Inches(0.6), Inches(3.4), Inches(12.2), Inches(0.4),
                "사전 압박 질문: 모순 1건당 1문항+ · 40점 미만 가치 드릴다운 · 개방형 질문 금지",
                size=11, color=MUTED)

    add_textbox(slide, Inches(0.6), Inches(4.0), Inches(12.2), Inches(0.4),
                "체인 ②  면접 그래프 (실시간, interview_graph.py)",
                size=13, bold=True, color=PRIMARY)

    chain2 = ["prepare_context", "analyze_answer",
              "generate_followups", "rank_and_filter"]
    cw2 = Inches(2.95)
    cy2 = Inches(4.45)
    for i, name in enumerate(chain2):
        x = cx + (cw2 + Inches(0.08)) * i
        add_arrow_step(slide, x, cy2, cw2, Inches(0.6), name,
                       accent=(name == "generate_followups"))

    bx = Inches(0.6)
    by = Inches(5.3)
    bw = Inches(6.0)
    bh = Inches(1.7)
    add_rect(slide, bx, by, bw, bh, fill=LIGHT, line=BORDER)
    add_textbox(slide, bx + Inches(0.2), by + Inches(0.1),
                bw - Inches(0.4), Inches(0.4),
                "question_source 분기", size=12, bold=True, color=INK)
    add_textbox(slide, bx + Inches(0.2), by + Inches(0.5),
                bw - Inches(0.4), bh - Inches(0.6),
                "pregenerated / custom / followup\n→ 꼬리질문의 꼬리질문도 생략하지 않도록\n   프롬프트에 출처를 강제 주입.",
                size=11, color=SUBINK, line_spacing=1.3)

    bx2 = Inches(6.85)
    add_rect(slide, bx2, by, bw, bh, fill=LIGHT, line=BORDER)
    add_textbox(slide, bx2 + Inches(0.2), by + Inches(0.1),
                bw - Inches(0.4), Inches(0.4),
                "asyncio.Queue 단일 워커", size=12, bold=True, color=INK)
    add_textbox(slide, bx2 + Inches(0.2), by + Inches(0.5),
                bw - Inches(0.4), bh - Inches(0.6),
                "services/followup_worker.py\n→ 동시 답변 제출에도 LLM 호출이 직렬화\n   순서 꼬임·중복 차단.",
                size=11, color=SUBINK, line_spacing=1.3)

    add_footer(slide, 9)


def slide_10(prs):
    """미션 2 [결과]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 2. 모순 탐지 + 꼬리질문  ·  [결과]",
              "시나리오별 입력 → 실제 질문 샘플")

    sy = Inches(1.95)
    sh = Inches(2.85)
    cols = [
        (DOCS / "모순.png", "모순 리스트"),
        (DOCS / "사전압박질문.png", "사전 압박 질문"),
        (DOCS / "면접.png", "실시간 면접 + 꼬리질문"),
    ]
    cw = Inches(4.05)
    cgap = Inches(0.12)
    sx = Inches(0.6)
    for i, (img, cap) in enumerate(cols):
        x = sx + (cw + cgap) * i
        add_rect(slide, x, sy, cw, sh, fill=LIGHT, line=BORDER)
        add_picture_fit(slide, img, x + Inches(0.05), sy + Inches(0.05),
                        cw - Inches(0.1), sh - Inches(0.1))
        add_textbox(slide, x, sy + sh + Inches(0.02), cw, Inches(0.3),
                    cap, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    qy = Inches(5.3)
    add_textbox(slide, Inches(0.6), qy, Inches(12.2), Inches(0.35),
                "질문 샘플 — 모든 질문에 basis (근거) · target_value (연결 가치) 필드 포함",
                size=12, bold=True, color=PRIMARY)
    samples = [
        '"이력서에는 99.1% 정확도를 750회 실험으로 달성했다고 되어 있는데, '
        '실험 1회당 평균 학습 시간과 GPU 사양을 알려주실 수 있나요?"',
        '"ONNX 가속으로 70% 단축이라 적혀 있는데, 단축 측정의 baseline이 무엇이고 '
        '어떤 디바이스에서 측정했나요?"',
        '"답변에 \"팀 전체가 같이 했다\"가 등장 → 자동 꼬리질문: '
        '\"본인이 단독 의사결정한 지점은 어디였나요?\""',
    ]
    add_rect(slide, Inches(0.6), qy + Inches(0.4),
             Inches(12.2), Inches(1.45),
             fill=WHITE, line=BORDER)
    add_bullets(slide, Inches(0.8), qy + Inches(0.5),
                Inches(11.8), Inches(1.3),
                samples, size=11, color=SUBINK, line_spacing=1.35)

    add_footer(slide, 10)


# ----- 확대 스크린샷 3장 (제목만) -----

def _slide_bigshot(prs, idx, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                title, size=26, bold=True, color=INK)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(1.0), Emu(36000),
             fill=PRIMARY, line=PRIMARY)

    bx = Inches(0.6)
    by = Inches(1.35)
    bw = Inches(12.2)
    bh = Inches(5.55)
    add_rect(slide, bx, by, bw, bh, fill=LIGHT, line=BORDER)
    add_picture_fit(slide, img_path,
                    bx + Inches(0.1), by + Inches(0.1),
                    bw - Inches(0.2), bh - Inches(0.2))
    add_footer(slide, idx)


def slide_11(prs):
    """모순 탐지 결과 — 확대 스크린샷."""
    _slide_bigshot(prs, 11, "모순 탐지 결과", DOCS / "모순.png")


def slide_12(prs):
    """사전 압박 질문 — 확대 스크린샷."""
    _slide_bigshot(prs, 12, "사전 압박 질문", DOCS / "사전압박질문.png")


def slide_13(prs):
    """실시간 면접 화면 — 확대 스크린샷."""
    _slide_bigshot(prs, 13, "실시간 면접 화면", DOCS / "면접.png")


def slide_14(prs):
    """미션 3 [수행]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 3. AX 대시보드  ·  [수행]",
              "후보자 등록 → 분석 → 리포트 → 면접 → 합격 결정까지 한 화면")

    lx = Inches(0.6)
    lw = Inches(5.5)
    add_textbox(slide, lx, Inches(1.95), lw, Inches(0.4),
                "아키텍처", size=14, bold=True, color=PRIMARY)

    arch = [
        ("백엔드", "FastAPI (port 8102)\n라우터 분리: auth · candidates · documents · analysis · interview"),
        ("프론트", "Next.js 14 App Router · TypeScript · Tailwind (port 3102)"),
        ("저장소", "SQLite + SQLAlchemy"),
        ("PDF", "reportlab 서버사이드 · 한글 폰트(Malgun Gothic) 자동 등록"),
    ]
    ay = Inches(2.4)
    ah = Inches(0.9)
    agap = Inches(0.08)
    for i, (h, b) in enumerate(arch):
        y = ay + (ah + agap) * i
        add_rect(slide, lx, y, lw, ah, fill=WHITE, line=BORDER)
        add_rect(slide, lx, y, Inches(1.1), ah, fill=INK, line=INK)
        add_textbox(slide, lx, y, Inches(1.1), ah,
                    h, size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, lx + Inches(1.25), y + Inches(0.08),
                    lw - Inches(1.4), ah - Inches(0.16),
                    b, size=10, color=SUBINK, line_spacing=1.25)

    # 좌하단: 안정성 박스 (신규)
    st_y = Inches(6.3)
    st_h = Inches(0.7)
    add_rect(slide, lx, st_y, lw, st_h, fill=PRIMARY, line=PRIMARY)
    add_textbox(slide, lx + Inches(0.2), st_y + Inches(0.07),
                lw - Inches(0.4), Inches(0.3),
                "안정성", size=11, bold=True, color=WHITE)
    add_textbox(slide, lx + Inches(0.2), st_y + Inches(0.32),
                lw - Inches(0.4), st_h - Inches(0.35),
                "call_claude_json 2회 재시도 · current_step 기반 재개 · API_SLEEP_SECONDS=0.5s로 rate limit 완충",
                size=9, color=WHITE, line_spacing=1.2)

    # 우측: 컴포넌트 맵
    rx = Inches(6.35)
    rw = Inches(6.5)
    add_textbox(slide, rx, Inches(1.95), rw, Inches(0.4),
                "핵심 컴포넌트 맵", size=14, bold=True, color=PRIMARY)
    add_rect(slide, rx, Inches(2.4), rw, Inches(4.6),
             fill=LIGHT, line=BORDER)

    comps = [
        ("ValueRadarChart", "12가치 레이더 시각화"),
        ("ValueList + EvidencePanel", "가치 클릭 → 원문 증거 슬라이드아웃"),
        ("ContradictionList", "모순 4축 + 심각도 색상 강조"),
        ("PreemptiveQuestions", "근거(basis)·연결 가치 함께 노출"),
        ("StatusSelector", "9단 상태(미분석~최종합격) 즉시 변경"),
        ("ProgressBar", "단계별 진행 + 예상 잔여 시간"),
        ("RawDataViewer", "OCR 원문 ↔ 구조화 JSON 탭 비교"),
        ("면접관 메모", "리포트 상단 1클릭 · 후보자 비노출"),
    ]
    cy = Inches(2.55)
    crow = Inches(0.52)
    for i, (n, d) in enumerate(comps):
        y = cy + crow * i
        add_textbox(slide, rx + Inches(0.25), y,
                    Inches(2.5), crow,
                    "•  " + n, size=11, bold=True, color=INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, rx + Inches(2.85), y,
                    rw - Inches(3.1), crow,
                    d, size=11, color=SUBINK,
                    anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 14)


def slide_15(prs):
    """미션 3 [결과] — 시연 영상 자리 (빈 슬라이드)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "미션 3. AX 대시보드  ·  [결과]",
              "동작하는 웹 프로토타입 — 5분 시연")

    # 중앙에 "시연 영상" 플레이스홀더 박스만
    px = Inches(2.0)
    py = Inches(2.3)
    pw = Inches(9.33)
    ph = Inches(4.4)
    add_rect(slide, px, py, pw, ph, fill=LIGHT, line=BORDER)
    add_textbox(slide, px, py, pw, ph,
                "▶  시연 영상",
                size=36, bold=True, color=MUTED,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 15)


def slide_16(prs):
    """엔지니어링 고민 ① 할루시네이션."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "엔지니어링 고민 ①  할루시네이션 자가검증 루프",
              "\"근거가 원문에 존재하는가\"를 기계가 스스로 확인하게 만들었습니다")

    lx = Inches(0.6)
    lw = Inches(5.5)
    ly = Inches(1.95)
    add_rect(slide, lx, ly, lw, Inches(1.7),
             fill=INK, line=INK)
    add_textbox(slide, lx + Inches(0.25), ly + Inches(0.15),
                lw - Inches(0.5), Inches(0.5),
                "문제 정의", size=14, bold=True, color=PRIMARY)
    add_textbox(slide, lx + Inches(0.25), ly + Inches(0.6),
                lw - Inches(0.5), Inches(1.05),
                "LLM이 examples 필드에 원문에 없는 문장을\n"
                "만들어 넣는 환각.\n"
                "→ 면접관이 리포트를 신뢰할 수 없는\n   치명적 실패 모드.",
                size=11, color=WHITE, line_spacing=1.35)

    dy = Inches(3.85)
    add_textbox(slide, lx, dy, lw, Inches(0.4),
                "검증 루프", size=12, bold=True, color=PRIMARY)
    diag_y = dy + Inches(0.5)
    nodes = [("score_12_values", False),
             ("self_verify_evidence", True),
             ("regenerate (필요 시)", False),
             ("doc_reliability", False)]
    nh = Inches(0.45)
    nw = Inches(2.55)
    for i, (n, accent) in enumerate(nodes):
        row = i // 2
        col = i % 2
        x = lx + (nw + Inches(0.15)) * col
        y = diag_y + (nh + Inches(0.18)) * row
        add_arrow_step(slide, x, y, nw, nh, n, accent=accent)

    rx = Inches(6.35)
    rw = Inches(6.5)
    add_textbox(slide, rx, Inches(1.95), rw, Inches(0.4),
                "4층 방어선", size=14, bold=True, color=PRIMARY)
    layers = [
        ("01  프롬프트 제약",
         "추상적 자기서술 불인정 · 원문 문장 그대로 인용 강제."),
        ("02  근거 필드 강제",
         "evidence · examples · basis · reasoning 전부 스키마 필수."),
        ("03  자가 점증 검증 루프 (핵심)",
         "substring + token-overlap ≥ 0.6 · 점수 40↑ & 검증 비율 < 50% →\n"
         "해당 가치만 1회 재호출해 인용 교체. verification 필드로 DB 영속화."),
        ("04  사람 검증 UI",
         "RawDataViewer — OCR 원문 ↔ 구조화 JSON 나란히 비교."),
    ]
    ly2 = Inches(2.4)
    lh = Inches(1.07)
    lgap = Inches(0.08)
    for i, (h, b) in enumerate(layers):
        y = ly2 + (lh + lgap) * i
        accent = "핵심" in h
        border = PRIMARY if accent else BORDER
        fill = LIGHT if accent else WHITE
        add_rect(slide, rx, y, rw, lh, fill=fill, line=border,
                 line_width=1.25 if accent else 0.75)
        add_textbox(slide, rx + Inches(0.2), y + Inches(0.1),
                    rw - Inches(0.4), Inches(0.4),
                    h, size=12, bold=True,
                    color=PRIMARY if accent else INK)
        add_textbox(slide, rx + Inches(0.2), y + Inches(0.45),
                    rw - Inches(0.4), lh - Inches(0.5),
                    b, size=10, color=SUBINK, line_spacing=1.3)

    add_footer(slide, 16)


def slide_17(prs):
    """엔지니어링 고민 ② 토큰·비용 실측 + Claude Code."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "엔지니어링 고민 ②  토큰·비용 실측 + Claude Code 활용",
              "대략치 메모가 아니라 DB에 남는 수치로 비용을 본다")

    # 좌측 상단: token_usage 테이블
    lx = Inches(0.6)
    lw = Inches(6.4)
    add_textbox(slide, lx, Inches(1.95), lw, Inches(0.4),
                "비용 계측 설계", size=14, bold=True, color=PRIMARY)

    tb_y = Inches(2.4)
    add_rect(slide, lx, tb_y, lw, Inches(1.45), fill=LIGHT, line=BORDER)
    add_textbox(slide, lx + Inches(0.2), tb_y + Inches(0.1),
                lw - Inches(0.4), Inches(0.35),
                "token_usage 테이블 · contextvars 자동 태깅",
                size=12, bold=True, color=INK)
    add_textbox(slide, lx + Inches(0.2), tb_y + Inches(0.45),
                lw - Inches(0.4), Inches(1.0),
                "candidate_id · phase · step · model · input/output tokens · cost_usd\n"
                "usage_scope(candidate_id, phase, step) 컨텍스트로\n"
                "LangGraph 노드 · OCR · Vision 호출이 자동 태깅되어 DB 기록.",
                size=10, color=SUBINK, line_spacing=1.35)

    # 좌측 하단: 실측 차트
    ch_y = Inches(3.95)
    ch_h = Inches(3.0)
    add_textbox(slide, lx, ch_y - Inches(0.02), lw, Inches(0.4),
                "실측 — step별 누적 비용 (USD, N=3 분석)",
                size=12, bold=True, color=PRIMARY)

    rows = fetch_step_costs()
    if rows:
        cats = [r[0] for r in rows]
        vals = [round(r[1], 4) for r in rows]
        chart_data = CategoryChartData()
        chart_data.categories = cats
        chart_data.add_series("누적 비용 (USD)", vals)
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            lx, ch_y + Inches(0.3), lw, ch_h - Inches(0.3),
            chart_data,
        )
        chart = gframe.chart
        chart.has_legend = False
        chart.has_title = False
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dls = plot.data_labels
            dls.font.size = Pt(9)
            dls.font.name = FONT
            dls.number_format = '"$"0.000'
            dls.position = XL_LABEL_POSITION.OUTSIDE_END
            # 주황색 막대
            ser = plot.series[0]
            fill = ser.format.fill
            fill.solid()
            fill.fore_color.rgb = PRIMARY
            ser.format.line.color.rgb = PRIMARY
        except Exception:
            pass
        # 축 폰트
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(9)
                axis.tick_labels.font.name = FONT
                axis.tick_labels.font.color.rgb = SUBINK
        except Exception:
            pass
    else:
        add_rect(slide, lx, ch_y + Inches(0.3), lw, ch_h - Inches(0.3),
                 fill=LIGHT, line=BORDER)
        add_textbox(slide, lx, ch_y + Inches(0.3), lw, ch_h - Inches(0.3),
                    "(DB 연결 실패 — 차트 생성 불가)",
                    size=11, color=MUTED,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 우측 상단: 측정값 요약 박스
    rx = Inches(7.2)
    rw = Inches(5.65)
    mv_y = Inches(1.95)
    mv_h = Inches(2.2)
    add_rect(slide, rx, mv_y, rw, mv_h, fill=PRIMARY, line=PRIMARY)
    add_textbox(slide, rx + Inches(0.25), mv_y + Inches(0.15),
                rw - Inches(0.5), Inches(0.4),
                "실측 요약 (N=3 분석, claude-sonnet-4-6)",
                size=11, bold=True, color=WHITE)
    add_textbox(slide, rx + Inches(0.25), mv_y + Inches(0.6),
                rw - Inches(0.5), Inches(0.7),
                "$0.38", size=36, bold=True, color=WHITE)
    add_textbox(slide, rx + Inches(0.25), mv_y + Inches(1.25),
                rw - Inches(0.5), Inches(0.3),
                "분석 1회 평균 · 분석 단계 누적 $1.14 / 19 호출",
                size=11, color=WHITE)
    add_textbox(slide, rx + Inches(0.25), mv_y + Inches(1.6),
                rw - Inches(0.5), Inches(0.4),
                "면접(8 호출) $0.19  ·  꼬리질문 1회 $0.03–0.05",
                size=11, color=WHITE)

    # 우측 하단: Claude Code 활용
    cc_y = Inches(4.25)
    cc_h = Inches(2.75)
    add_rect(slide, rx, cc_y, rw, cc_h, fill=WHITE, line=BORDER)
    add_textbox(slide, rx + Inches(0.25), cc_y + Inches(0.15),
                rw - Inches(0.5), Inches(0.4),
                "최신 AI 도구 활용 — Claude Code",
                size=13, bold=True, color=PRIMARY)
    add_textbox(slide, rx + Inches(0.25), cc_y + Inches(0.55),
                rw - Inches(0.5), Inches(0.45),
                "claude-opus-4-7 을 프로그래밍 모델로 사용 (서비스 자체는 sonnet-4-6)",
                size=10, color=SUBINK, line_spacing=1.3)
    cc_bullets = [
        "LangGraph 노드 분리 설계 검토",
        "자가검증 루프 임계값(token-overlap 0.6) 결정 토론",
        "reportlab 한글 폰트 등록 트러블슈팅",
        "본 발표 자료(스크립트·차트 명세) 초안 생성",
    ]
    add_bullets(slide, rx + Inches(0.25), cc_y + Inches(1.05),
                rw - Inches(0.5), Inches(1.3),
                cc_bullets, size=10, color=SUBINK, line_spacing=1.3)
    add_textbox(slide, rx + Inches(0.25), cc_y + cc_h - Inches(0.4),
                rw - Inches(0.5), Inches(0.35),
                "원칙 — 생성된 코드는 100% 직접 읽고 합리화한 후 채택.",
                size=10, bold=True, color=INK)

    add_footer(slide, 17)


def slide_18(prs):
    """솔직히 못 한 부분 및 개선하고 싶은 부분."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "솔직히 못 한 부분 및 개선하고 싶은 부분",
              "보강해야 할 4가지 — 어떤 순서로 풀 것인지")

    lx = Inches(0.6)
    lw = Inches(7.0)
    add_textbox(slide, lx, Inches(1.85), lw, Inches(0.4),
                "미흡 항목과 그 이유", size=14, bold=True, color=PRIMARY)

    items = [
        ("①  프롬프트 회귀 테스트 부재",
         "고정 픽스처 × 계약 검증을 pytest 스위트로 만들지 못함. MVP 우선순위에서 밀림."),
        ("②  Prompt Cache 미적용",
         "시스템 프롬프트·12가치 정의를 Anthropic prompt cache로 고정 못함. 토큰 절감 여지."),
        ("③  다중 평가자 협업·권한 모델 미구현",
         "단일 사용자 시나리오에 집중. 평가자별 권한·교차 평가·합산 합격 흐름 미구현."),
        ("④  LLM 평가 자동화(Eval) 부재",
         "Ragas/LangSmith 같은 도구로 모순·꼬리질문 품질을 정량화 못함. 사람 검수에 의존."),
    ]
    iy = Inches(2.3)
    ih = Inches(1.1)
    igap = Inches(0.08)
    for i, (h, b) in enumerate(items):
        y = iy + (ih + igap) * i
        add_rect(slide, lx, y, lw, ih, fill=WHITE, line=BORDER)
        add_rect(slide, lx, y, Inches(0.12), ih, fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, lx + Inches(0.25), y + Inches(0.13),
                    lw - Inches(0.4), Inches(0.4),
                    h, size=12, bold=True, color=INK)
        add_textbox(slide, lx + Inches(0.25), y + Inches(0.5),
                    lw - Inches(0.4), ih - Inches(0.55),
                    b, size=10, color=SUBINK, line_spacing=1.3)

    # 우측: 개선사항 1/2/3
    rx = Inches(7.85)
    rw = Inches(5.0)
    add_textbox(slide, rx, Inches(1.85), rw, Inches(0.4),
                "개선하고 싶은 부분", size=14, bold=True, color=PRIMARY)

    improvements = [
        ("개선사항 1",
         "①  해소.\n고정 이력서 픽스처 5종 × 계약 검증 pytest 스위트 + CI."),
        ("개선사항 2",
         "②  해소.\nPrompt cache를 시스템·12가치·OCR 지침에 적용 →\n토큰 비용 약 30% 추가 절감."),
        ("개선사항 3",
         "④  첫 걸음.\n모순·꼬리질문 회귀용 골든 데이터셋 30건 +\n자동 평가 파이프라인 초안."),
    ]
    my = Inches(2.3)
    mh = Inches(1.5)
    mgap = Inches(0.1)
    for i, (h, b) in enumerate(improvements):
        y = my + (mh + mgap) * i
        add_rect(slide, rx, y, rw, mh, fill=PRIMARY, line=PRIMARY)
        add_textbox(slide, rx + Inches(0.25), y + Inches(0.12),
                    rw - Inches(0.5), Inches(0.4),
                    h, size=14, bold=True, color=WHITE)
        add_textbox(slide, rx + Inches(0.25), y + Inches(0.5),
                    rw - Inches(0.5), mh - Inches(0.55),
                    b, size=10, color=WHITE, line_spacing=1.3)

    add_textbox(slide, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.35),
                "무엇을 못 했는지 알고 있다는 것 + 어떤 순서로 풀 것인지가 정해져 있다는 것.",
                size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 18)


def slide_19(prs):
    """클로징."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5),
             fill=PRIMARY, line=PRIMARY)
    add_rect(slide, Inches(11.6), Inches(0.6), Inches(1.4), Inches(1.4),
             fill=PRIMARY, line=PRIMARY)

    add_textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.6),
                "한 줄 메시지", size=16, color=MUTED)

    add_textbox(slide, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.0),
                "돌아가는 AI가 아니라,",
                size=44, bold=True, color=INK)
    add_textbox(slide, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.0),
                "분기마다 더 싸지고 더 정확해지는 AI를 만들겠습니다.",
                size=36, bold=True, color=PRIMARY)

    add_textbox(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.5),
                "·  0.1%의 효율을 측정 가능한 인프라 위에서 쌓아가겠습니다.",
                size=14, color=SUBINK)
    add_textbox(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.5),
                "·  인터엑스의 12가치를 회귀 테스트·비용 계측·자가 검증 루프로 옮기겠습니다.",
                size=14, color=SUBINK)

    add_rect(slide, Inches(0.9), Inches(6.2), Inches(11.5), Emu(18000),
             fill=BORDER, line=BORDER)
    add_textbox(slide, Inches(0.9), Inches(6.4), Inches(7), Inches(0.4),
                "감사합니다.  질문 받겠습니다.", size=16, bold=True, color=INK)
    add_textbox(slide, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.3),
                "주재홍   ·   kinduki123@gmail.com   ·   github.com/flyingjoojak",
                size=11, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    slide_17(prs)
    slide_18(prs)
    slide_19(prs)

    prs.save(str(OUT))
    print(f"saved: {OUT}  (slides={len(prs.slides)})")


if __name__ == "__main__":
    main()

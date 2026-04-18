from pptx import Presentation


def extract_ppt_text(file_path: str) -> str:
    prs = Presentation(file_path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text = shape.text_frame.text
                if slide_text:
                    texts.append(slide_text)
    return "\n".join(texts)
